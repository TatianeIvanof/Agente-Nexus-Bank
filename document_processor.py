import os
import os
import pandas as pd
import json
from dataclasses import dataclass
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation

# Estrutura simples e nativa para substituir o LCDocument do LangChain
@dataclass
class SimpleDocument:
    page_content: str
    metadata: dict

def load_documents_from_folder(folder_path="docs"):
    documents = []
    if not os.path.exists(folder_path):
        os.makedirs(folder_path)
        return documents

    for file_name in os.listdir(folder_path):
        file_path = os.path.join(folder_path, file_name)
        ext = file_name.split('.')[-1].lower()
        content = ""

        try:
            if ext == "pdf":
                from pypdf import PdfReader
                reader = PdfReader(file_path)
                for i, page in enumerate(reader.pages):
                    text = page.extract_text() or ""
                    if text.strip():
                        documents.append(SimpleDocument(
                            page_content=text,
                            metadata={"source": file_name, "page": i + 1, "type": "PDF"}
                        ))
                continue

            elif ext == "docx":
                doc = Document(file_path)
                content = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

            elif ext in ["xlsx", "xls"]:
                excel_data = pd.read_excel(file_path, sheet_name=None)
                for sheet_name, df in excel_data.items():
                    content += f"\nAba: {sheet_name}\n" + df.to_string()

            elif ext == "pptx":
                prs = Presentation(file_path)
                for i, slide in enumerate(prs.slides):
                    slide_text = ""
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text += shape.text + "\n"
                    if slide_text.strip():
                        documents.append(SimpleDocument(
                            page_content=slide_text,
                            metadata={"source": file_name, "slide": i + 1, "type": "PowerPoint"}
                        ))
                continue

            elif ext in ["md", "txt"]:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()

            elif ext == "csv":
                df = pd.read_csv(file_path)
                content = df.to_string()

            elif ext == "json":
                with open(file_path, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    content = json.dumps(data, ensure_ascii=False, indent=2)

            elif ext == "html":
                with open(file_path, "r", encoding="utf-8") as f:
                    soup = BeautifulSoup(f.read(), "html.parser")
                    content = soup.get_text(separator="\n")

            if content.strip():
                documents.append(SimpleDocument(
                    page_content=content,
                    metadata={"source": file_name, "type": ext.upper()}
                ))

        except Exception as e:
            print(f"Erro ao ler o arquivo {file_name}: {e}")

    return documents

# Função nativa para divisão de texto (chunking) sem dependências externas
def split_documents(documents, chunk_size=800, chunk_overlap=150):
    chunks = []
    for doc in documents:
        text = doc.page_content
        start = 0
        while start < len(text):
            end = start + chunk_size
            chunk_text = text[start:end]
            chunks.append(SimpleDocument(
                page_content=chunk_text,
                metadata=doc.metadata
            ))
            start += chunk_size - chunk_overlap
    return chunks